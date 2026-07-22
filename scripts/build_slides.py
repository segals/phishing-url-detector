"""Build the ~15-minute slide deck (presentation/url_detector.pptx) from the results.

Every slide carries speaker notes (the fuller explanation to say out loud), so the
deck is self-contained for presenting. Course concepts are introduced and explained
where they are first used, and flagged [COURSE] in the notes.
"""
import os, pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

INK = RGBColor(0x1d, 0x2b, 0x3a); ACCENT = RGBColor(0x2a, 0x9d, 0x8f)
RED = RGBColor(0xe7, 0x6f, 0x51); GREY = RGBColor(0x5b, 0x64, 0x72); WHITE = RGBColor(0xff, 0xff, 0xff)
GOLD = RGBColor(0xb8, 0x86, 0x2e)
TAB, FIG = "results/tables", "results/figures"

def val(tab, col, contains, valcol):
    d = pd.read_csv(f"{TAB}/{tab}.csv"); r = d[d[col].astype(str).str.contains(contains)]
    return str(r.iloc[0][valcol]) if len(r) else "?"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def _box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tb.text_frame.word_wrap = True
    return tb.text_frame

def _run(p, text, size, color=INK, bold=False, italic=False):
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic; f.name = "Calibri"

def header(s, title, kicker=None, tag=None):
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tf = _box(s, 0.6, 0.4, 12.1, 1.25)
    p = tf.paragraphs[0]; _run(p, title, 29, INK, bold=True)
    if kicker:
        p2 = tf.add_paragraph(); _run(p2, kicker, 14, ACCENT, italic=True)
    if tag:
        tb = _box(s, 10.7, 0.42, 2.5, 0.5); tb.paragraphs[0].alignment = 2
        _run(tb.paragraphs[0], tag, 13, RED, bold=True)

def bullets(s, items, top=1.95, size=17.5, left=0.7, width=12.0):
    tf = _box(s, left, top, width, 5.2)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        lvl = 0
        if it.startswith("  "): lvl = 1; it = it.strip()
        p.level = lvl
        color = INK if lvl == 0 else GREY
        _run(p, ("•  " if lvl == 0 else "–  ") + it, size - lvl*1.5, color)

def picture(s, name, left=7.0, top=1.95, width=6.0):
    p = f"{FIG}/{name}.png"
    if os.path.exists(p): s.shapes.add_picture(p, Inches(left), Inches(top), width=Inches(width))

def notes(s, text): s.notes_slide.notes_text_frame.text = text
def slide(): return prs.slides.add_slide(BLANK)

def bottom_note(s, text, color=RED):
    tf = _box(s, 0.7, 6.55, 12, 0.7); _run(tf.paragraphs[0], text, 17, color, bold=True)

# =====================================================================
# 1 — TITLE
# =====================================================================
s = slide()
bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background()
tf = _box(s, 0.9, 2.3, 11.5, 3.2)
_run(tf.paragraphs[0], "Catching phishing from the link alone", 40, WHITE, bold=True)
p = tf.add_paragraph(); _run(p, "…and why “99% accuracy” in URL phishing detection is mostly an illusion", 20, RGBColor(0x9a,0xd1,0xc9), italic=True)
p = tf.add_paragraph(); p.space_before = Pt(26)
_run(p, "Data Science & Cyber  ·  the URL / link detector (Student B)", 15, RGBColor(0xc7,0xd0,0xd8))
notes(s, "~15 minute talk. Opening line: 'Can you tell if a link is phishing just from the URL string "
         "itself — no clicking, no page content?' I built a detector that does exactly that, got a "
         "near-perfect score, and then spent most of the project proving that score doesn't mean what "
         "it looks like. That honest audit is the real contribution.")

# =====================================================================
# 2 — THE QUESTION & THE PROJECT
# =====================================================================
s = slide(); header(s, "The question, and the project")
bullets(s, ["Phishing is the #1 way attackers get a first foothold — and much of it is just a link.",
            "Can we classify a URL as phishing / legitimate from the STRING ALONE?",
            "  no page fetch → fast (~milliseconds), private, works before the user clicks",
            "Two-part project: this half = the URL / link channel; my partner = the email content channel.",
            "Built independently. Plan: build a detector → then stress-test whether its score is real."])
notes(s, "Frame the two-part project: my partner reads the words of the email; I read the link. Both are "
         "meant to be fused later (out of scope here). The plan is deliberately 'build then break' — the "
         "20% that is model-building sets up the 80% that is critical evaluation.")

# =====================================================================
# 3 — THE DATA
# =====================================================================
s = slide(); header(s, "The data", "real, public, recent")
bullets(s, ["PhiUSIIL (Prasad & Chandra, 2024; UCI id 967) — 235,795 real URLs, 43% phishing",
            "Preprocessing: de-duplicate on the URL BEFORE splitting (no leakage), keep the natural",
            "  class balance (no rebalancing), stratified 60/20/20 split, seed 42, test used once",
            "For the generalisation test later: an INDEPENDENT Kaggle set (651k URLs) + live OpenPhish",
            "Reproducible: cleaned data is SHA-256 hashed and logged"], width=6.4)
picture(s, "F1_eda", left=6.95, top=2.15, width=6.15)
notes(s, "[COURSE: EDA + train/val/test methodology.] Data prep was minimal on purpose: only a label "
         "standardisation (1=phishing), dropping empties, and de-duplicating identical URLs before the "
         "split so the same URL can't appear in train and test. I did NOT rebalance classes, remove "
         "outliers, or scale the raw data — so the near-perfect score I'm about to show is a property of "
         "the RAW data, not of my cleaning. Right figure: the EDA already looks suspicious — the classes "
         "separate almost perfectly on trivial quantities, which never happens with genuinely hard data.")

# =====================================================================
# 4 — HOW WE PICKED THE FEATURES
# =====================================================================
s = slide(); header(s, "How we chose the features", "feature engineering")
bullets(s, ["A URL is a short, STRUCTURED string — not prose.",
            "So unlike my partner (who vectorises text with TF-IDF / bag-of-words), I compute",
            "  interpretable numbers DIRECTLY from the URL — no learned text embedding.",
            "Why interpretable, hand-built features (not a deep character model / URLNet)?",
            "  → so I can later attribute the score to specific, human-readable causes",
            "40 features, grouped into 6 families (next slide)"])
notes(s, "[COURSE: feature engineering.] This is the key design choice. Text has to be turned into "
         "numbers somehow — my partner uses TF-IDF (term frequency × inverse document frequency) over "
         "words. A URL isn't a sentence; it's structured (scheme://host/path?query). So I hand-engineer "
         "40 numeric features straight from the string. I deliberately avoided a black-box deep model "
         "(like URLNet) because interpretability is the whole point: when I find the near-perfect score "
         "is an artifact, I want to name exactly which feature is responsible — which a black box would "
         "hide.")

# =====================================================================
# 5 — THE SIX FEATURE FAMILIES
# =====================================================================
s = slide(); header(s, "The 40 features — six families")
bullets(s, ["Lexical / length — URL length; counts of dots, slashes, digits, hyphens; digit ratio",
            "Host / domain — number of subdomains; is the host a raw IP; unusual port",
            "TLD — the suffix (.com vs .tk vs .xyz); is it an abuse-associated TLD",
            "Entropy — randomness of the host string (algorithmically-generated domains)",
            "Scheme / encoding — HTTPS?; punycode (xn--); percent/hex encoding",
            "Brand / cue — words like ‘verify/login/secure’; look-alike distance to real brands"],
        size=17)
notes(s, "Walk the six families briefly. Lexical = size and punctuation of the string. Host/domain = "
         "structure of the hostname (an IP literal instead of a domain name is suspicious). TLD = the "
         "suffix; some cheap TLDs (.tk, .xyz) are over-represented in abuse feeds. Entropy and Brand "
         "each have real theory — next slide. Scheme/encoding includes HTTPS (which turns out to be the "
         "villain of the story) and punycode, the xn-- encoding used to register look-alike "
         "international-character domains.")

# =====================================================================
# 6 — TWO FEATURES WITH REAL THEORY
# =====================================================================
s = slide(); header(s, "Two features worth the math", "information theory · edit distance")
tf = _box(s, 0.7, 1.9, 12, 4.4)
_run(tf.paragraphs[0], "Shannon entropy of the host:   H = − Σ p · log₂(p)", 22, INK, bold=True)
p = tf.add_paragraph(); p.space_after = Pt(6)
_run(p, "“google” = repeating, predictable letters → LOW entropy;  “x7qz9wvt” ≈ random → HIGH entropy.", 16, GREY)
p = tf.add_paragraph(); _run(p, "→ catches algorithmically-generated (DGA) phishing domains.", 16, ACCENT, italic=True)
p = tf.add_paragraph(); p.space_before = Pt(20)
_run(p, "Brand look-alike:   Levenshtein edit distance to a known brand", 22, INK, bold=True)
p = tf.add_paragraph(); p.space_after = Pt(6)
_run(p, "edit distance = min single-character inserts / deletes / substitutions to turn one word into another.", 16, GREY)
p = tf.add_paragraph(); _run(p, "“paypa1” is edit-distance 1 from “paypal”  → flag distance 1–2 as a typosquat.", 16, ACCENT, italic=True)
notes(s, "[COURSE: Shannon entropy / information theory; edit distance.] Entropy formula: sum over the "
         "character distribution of the host, p·log2(p). Intuitively it measures unpredictability — a "
         "real brand name is low-entropy, a random string is high-entropy. Levenshtein distance is the "
         "classic string-similarity measure, computed by dynamic programming; paypa1→paypal is one "
         "substitution, so distance 1. Domains within edit distance 1–2 of a known brand get flagged as "
         "typosquats. These two features are where the URL side borrows the same math toolkit as the "
         "rest of the course.")

# =====================================================================
# 7 — IT LOOKS PERFECT
# =====================================================================
s = slide(); header(s, "Step 1 — train it. It looks perfect: F1 ≈ 1.0")
bullets(s, [f"All provided features → F1 { val('T0_artifact_decomposition','representation','provided','f1') }",
            f"ONE feature (URLSimilarityIndex) alone → F1 { val('T0_artifact_decomposition','representation','URLSimilarityIndex','f1') }  → basically the label (leakage)",
            f"My honest hand-built lexical features → F1 { val('T0_artifact_decomposition','representation','honest','f1') }",
            "Accuracy paradox: always-predict-legit = 57% accuracy but 0% recall (catches nothing)",
            "A one-line ‘not-HTTPS → phishing’ rule alone already scores F1 0.68"], width=12)
bottom_note(s, "This is where the FIRST learning ends — and the real project begins.")
notes(s, "[COURSE: the accuracy paradox.] This is the model you 'got' in the first round. It's excellent "
         "— and that's the problem. The decomposition is the first alarm: a single provided column "
         "scores 0.995, meaning it's essentially a copy of the answer key (leakage). And the accuracy "
         "paradox: 57% of the data is legit, so a do-nothing model looks 57% accurate while catching "
         "zero phishing — which is why I never report accuracy alone. Everything after this slide is a "
         "'discovery' about why the perfect score is fake.")

# =====================================================================
# 8 — DISCOVERY 1: THE ARTIFACT
# =====================================================================
s = slide(); header(s, "Discovery 1 — the score is a collection artifact", tag="DISCOVERY 1")
bullets(s, [f"Legitimate URLs in PhiUSIIL: { val('T3a_https_artifact','class','legitimate','https_rate') } HTTPS.   Phishing: { val('T3a_https_artifact','class','phishing','https_rate') } HTTPS.",
            "So the model largely learns ‘HTTP = phishing’ — a fact about how the data was collected,",
            "  NOT about phishing (real modern phishing is mostly HTTPS — this rule is backwards!)",
            "Legit URLs here are clean home-pages; phishing are messy deep links.",
            "The model learned the dataset’s shape, not danger."], width=12)
notes(s, "Discovery 1: WHY is it perfect? Legit URLs are 100% HTTPS, phishing only 49% — a 51-point gap "
         "baked into how the two classes were gathered. The model is basically an HTTPS detector. This "
         "is the core critical finding: near-perfect benchmark scores measure dataset construction, not "
         "the phishing problem. Note the irony — in the real world phishing is mostly HTTPS now, so this "
         "learned rule is actually backwards outside the dataset.")

# =====================================================================
# 9 — HOW WE EVALUATE HONESTLY (concepts)
# =====================================================================
s = slide(); header(s, "How we evaluate honestly", "the metric toolkit — and why")
bullets(s, ["Accuracy lies under imbalance (paradox) → report Precision, Recall, F1",
            "F2 — like F1 but weights RECALL more: a missed phish (FN) costs more than a false alarm (FP)",
            "MCC (Matthews) — one score in [−1,1] using ALL 4 confusion cells; honest under imbalance",
            "AUC — P(a random phish scores above a random legit); 0.5 = coin-flip",
            "Bootstrap 95% CI — resample the test set 1000× to get error bars on every metric",
            "McNemar test — is model A really better than B, or just luck on this split?"], size=17)
notes(s, "[COURSE: evaluation metrics + statistical validity.] Explain the two that get asked about. "
         "MCC = Matthews Correlation Coefficient: computed from true/false positives and negatives all "
         "four together, it ranges −1 (perfectly wrong) to +1 (perfect), 0 = random. Unlike accuracy or "
         "even F1, you can't fake a good MCC by exploiting the majority class — that's why it's the "
         "honest single number under imbalance. Bootstrap CI: instead of trusting one number from one "
         "test set, I resample the test set with replacement 1000 times, recompute the metric each time, "
         "and report the middle 95% — it tells me how much the number would wobble, so I don't over-read "
         "a 0.001 difference. McNemar does the same job for comparing two models.")

# =====================================================================
# 10 — DISCOVERY 2: SHAP
# =====================================================================
s = slide(); header(s, "Discovery 2 — SHAP confirms it: is_https drives everything", tag="DISCOVERY 2")
bullets(s, ["SHAP = a faithful way to credit each feature for the prediction (game theory)",
            f"#1 driver = is_https, ~3× the next feature",
            "The model is a thin veneer over the HTTPS artifact —",
            "  which is exactly why the attacks on the next slide work"], width=6.4)
picture(s, "F5_shap", left=7.0, top=2.1, width=6.0)
notes(s, "[COURSE: explainability / SHAP.] SHAP (Lundberg & Lee 2017) assigns each feature a fair share "
         "of credit for a prediction, borrowing the Shapley-value idea from game theory — it's faithful, "
         "not a guess. Result: is_https dominates by ~3×. So Discovery 1 (the HTTPS artifact) and "
         "Discovery 2 (SHAP) independently point at the same feature. That sets up the attacks: if the "
         "model leans on HTTPS, an attacker just needs to flip that.")

# =====================================================================
# 11 — DISCOVERY 3: ARMS RACE
# =====================================================================
s = slide(); header(s, "Discovery 3 — attack it: robust to tricks, but not to the artifact", tag="DISCOVERY 3")
bullets(s, [f"Homoglyph / typosquat (char tricks) → recall stays ~0.99",
            "  a STRUCTURAL detector ignores letters → resists what BREAKS a text model (nice contrast!)",
            f"Just switch HTTP → HTTPS → recall { val('T3_arms_race','attack','https_upgrade','recall_attacked') }",
            f"Host on a clean HTTPS domain (mimicry) → recall { val('T3_arms_race','attack','homepage_mimicry','recall_attacked') }  (total evasion)",
            "Normalisation defeats homoglyphs; NOTHING URL-only defeats mimicry"], width=6.5)
picture(s, "F3_arms_race", left=7.25, top=2.3, width=5.85)
notes(s, "[COURSE: adversarial ML — evasion, perturbation, homoglyphs, adversarial training.] The "
         "attacker keeps the URL working but reshapes its surface. Two findings: (1) character-level "
         "attacks (Cyrillic look-alikes, typosquats) barely move recall — a structural detector doesn't "
         "care about the exact letters, the opposite of a bag-of-words text model, which those same "
         "tricks defeat. Nice complementarity with my partner. (2) But the artifact is the hole: just "
         "serving over HTTPS cuts recall to 0.67, and a clean bare HTTPS domain evades completely "
         "(recall 0). This is what real phishing does — host on legitimate cloud infrastructure.")

# =====================================================================
# 12 — UNSUPERVISED CHECK
# =====================================================================
s = slide(); header(s, "Cross-check — can we catch it with NO labels?", "anomaly detection")
bullets(s, ["Train only on LEGITIMATE URLs; flag anything unusual as phishing",
            f"Isolation Forest → AUC { val('T10_anomaly','unsupervised_detector','IsolationForest','auc') }   (isolates odd points with random splits)",
            f"Local Outlier Factor → AUC { val('T10_anomaly','unsupervised_detector','LocalOutlierFactor','auc') }   (compares a point’s density to its neighbours’)",
            "Respectable with zero phishing labels — but below supervised, and riding the same artifacts"])
notes(s, "[COURSE: abnormality / anomaly detection — Isolation Forest, LOF.] A weaker-assumption sanity "
         "check: with no phishing labels at all, unsupervised detectors trained only on legit URLs still "
         "get AUC 0.83 / 0.94. Isolation Forest isolates anomalies quickly with random splits; LOF flags "
         "points in low-density regions relative to their neighbours. Useful, but still built on the same "
         "structural signals, so they inherit the same fragility.")

# =====================================================================
# 13 — DISCOVERY 4: CROSS-DATASET
# =====================================================================
s = slide(); header(s, "Discovery 4 — it does NOT generalise", tag="DISCOVERY 4")
bullets(s, [f"Train PhiUSIIL, test an INDEPENDENT 651k dataset: F1 { val('T4_cross_dataset','setting','in-distribution','f1') } → { val('T4_cross_dataset','setting','full model','f1') }",
            f"AUC { val('T4_cross_dataset','setting','in-distribution','auc') } → { val('T4_cross_dataset','setting','full model','auc') }  — below 0.5 = ‘confidently backwards’",
            "  it doesn’t just lose skill; its ranking INVERTS (legit scored above phishing)",
            "  both classes collapse to ~1.0 → the model can’t tell them apart at all",
            f"BUT live OpenPhish phishing: recall { val('T4_cross_dataset','setting','OpenPhish','recall_phishing') } (today’s phishing is still messy → caught)"], width=6.5)
picture(s, "F4_cross_dataset", left=7.25, top=2.35, width=5.85)
notes(s, "[COURSE: domain shift / concept drift; AUC.] The real test — a totally different dataset. The "
         "near-perfect model collapses to F1 0.67, AUC 0.43. AUC below 0.5 is the key idea: AUC is the "
         "chance a random phishing URL scores above a random legit one; 0.5 is a coin flip. Below 0.5 "
         "means the model is systematically ranking legit ABOVE phishing — it learned a real pattern but "
         "that pattern is inverted in the new data. I call this 'confidently backwards': not clueless, "
         "but confidently wrong. The boxplot shows why: in-distribution the classes separate cleanly; on "
         "the new data both pile up at ~1.0. Yet it still catches 100% of live current phishing, because "
         "real phishing today is still structurally messy — a genuinely nuanced result.")

# =====================================================================
# 14 — WHY IT INVERTS (distance metrics)
# =====================================================================
s = slide(); header(s, "Why it inverts — distance metrics pinpoint it", "KS & Wasserstein")
bullets(s, ["Measure how far each feature’s distribution SHIFTS between the two datasets",
            "  KS statistic = biggest gap between the two cumulative distributions",
            "  Wasserstein (earth-mover) = work to turn one distribution into the other",
            "Most-shifted features = is_https, path_len, n_subdomain",
            "  …which are EXACTLY the top SHAP drivers → it built on the least transferable signals"], width=6.5)
picture(s, "F7_domain_shift", left=7.25, top=2.3, width=5.85)
notes(s, "[COURSE: distance metrics — Kolmogorov–Smirnov, Wasserstein/EMD.] This closes the loop. I "
         "measure, feature by feature, how different the two datasets are, using two standard distances: "
         "KS (largest gap between the two cumulative distribution curves) and Wasserstein (the earth-"
         "mover distance — how much probability mass must move, and how far). The most-shifted features "
         "turn out to be precisely is_https, path_len, n_subdomain — the same features SHAP flagged as "
         "the model's top drivers. So the collapse isn't mysterious: the model built its decision on the "
         "three least-transferable signals. Two different tools (SHAP + distance metrics) corroborate "
         "one causal story.")

# =====================================================================
# 15 — CONCLUSION
# =====================================================================
s = slide(); header(s, "Conclusion — a useful second opinion, not a solution")
bullets(s, ["URL structure alone is NOT a generalisable phishing signal:",
            "  built on a dataset artifact (HTTPS + correlated structure) → trivially evaded, collapses OOD",
            "But it IS fast, private, interpretable, and robust to character obfuscation",
            "Right design = FUSION: this link channel + my partner’s content channel + reputation/age",
            "  → beat one channel and the other still fires"])
bottom_note(s, "Near-perfect URL benchmarks measure dataset construction, not phishing.")
notes(s, "The honest bottom line. The detector isn't worthless — it's fast, private, interpretable, and "
         "resists cheap character tricks, and it catches today's messy phishing. But it can't stand "
         "alone: it's an artifact-driven benchmark score that an attacker beats by using legitimate "
         "hosting, and it doesn't survive a change of dataset. Hence fusion with the content channel — "
         "the two look at completely independent evidence.")

# =====================================================================
# 16 — COURSE CONCEPTS MAP
# =====================================================================
s = slide(); header(s, "Course concepts used in this project")
bullets(s, ["EDA · distributions · skewness/kurtosis · robust stats · Pearson vs Spearman · Mann–Whitney U",
            "Feature engineering · Shannon entropy (information theory) · Levenshtein edit distance",
            "Accuracy paradox · Precision/Recall/F1/F2/MCC/AUC · cost-sensitive thresholds",
            "Bootstrap confidence intervals · McNemar’s test",
            "Explainability (SHAP) · adversarial ML (evasion, homoglyphs, adversarial training)",
            "Anomaly detection (Isolation Forest, LOF) · calibration (Brier)",
            "Domain shift / concept drift · distance metrics (KS, Wasserstein) · feature ablation · error analysis"],
        size=15.5)
notes(s, "The explicit concept map — this is the slide to point to if the lecturer asks 'where did you "
         "use what we taught?'. Every item here is earned in a specific experiment, and the mapping is "
         "also in the report appendix and the notebook's final section.")

# =====================================================================
# 17 — THANKS
# =====================================================================
s = slide()
bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background()
tf = _box(s, 0.9, 2.9, 11.5, 2.2)
_run(tf.paragraphs[0], "Thank you", 40, WHITE, bold=True)
p = tf.add_paragraph(); _run(p, "github.com/segals/phishing-url-detector   ·   reproducible: python run_all.py", 16, RGBColor(0x9a,0xd1,0xc9))
p = tf.add_paragraph(); p.space_before = Pt(10)
_run(p, "Questions?", 18, RGBColor(0xc7,0xd0,0xd8), italic=True)
notes(s, "Likely questions: (1) 'Why is your score so high?' → it's an artifact, that's the whole point. "
         "(2) 'Did you cause it by preprocessing?' → no; only dedup + split, no rebalancing, hash logged. "
         "(3) 'Biggest weakness?' → doesn't generalise (AUC 0.43) and evaded by legit hosting → fusion. "
         "(4) 'Why not deep learning / URLNet?' → interpretability, so I can name the artifact.")

os.makedirs("presentation", exist_ok=True)
prs.save("presentation/url_detector.pptx")
print(f"wrote presentation/url_detector.pptx ({len(prs.slides._sldIdLst)} slides, with speaker notes)")
